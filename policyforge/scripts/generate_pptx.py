#!/usr/bin/env python3
"""
Generate PolicyForge PowerPoint Presentation (.pptx)

Requirements (from 10-slide outline):
1. Title Slide
2. The $30B Problem Statement
3. The Solution: PolicyForge Overview
4. System Architecture & Tech Stack
5. LLM Orchestration Depth & The Grounding Gate
6. Data Engineering & Analytics (DuckDB Pipeline)
7. Operational Results & Metrics
8. Managing Clinical Risk & Human-in-the-Loop
9. Strategic Recommendations for Cotiviti
10. Conclusion & Core Takeaways

Design philosophy:
- Crisp white/light gray backgrounds (#F8F9FA)
- Deep Navy/Charcoal text (#1E293B)
- Vibrant Cyan/Teal accents (#0EA5E9)
- Huge metrics, no walls of text
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# ── Brand palette ─────────────────────────────────────────────────────────────
NAVY       = RGBColor(0x1E, 0x29, 0x3B)  # Deep Navy/Charcoal
CYAN       = RGBColor(0x0E, 0xA5, 0xE9)  # Vibrant Cyan/Teal Accent
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF8, 0xF9, 0xFA)  # Crisp Light Gray
MID_GRAY   = RGBColor(0x64, 0x74, 0x8B)  # For muted labels
RED        = RGBColor(0xEF, 0x44, 0x44)
GREEN      = RGBColor(0x10, 0xB9, 0x81)
AMBER      = RGBColor(0xF5, 0x9E, 0x0B)

SLIDE_W = Inches(13.33)   # widescreen 16:9
SLIDE_H = Inches(7.5)


# ─────────────────────────────────────────────────────────────────────────────
# Low-level helpers
# ─────────────────────────────────────────────────────────────────────────────

def add_rect(slide, left, top, width, height, fill_color, line_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()
    return shape

def add_chevron(slide, left, top, width, height, fill_color):
    shape = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape

def tf_para(tf, text, size, bold=False, color=NAVY, align=PP_ALIGN.LEFT, space_before=0):
    para = tf.add_paragraph()
    para.text = text
    para.alignment = align
    run = para.runs[0]
    # Set to a highly readable standard sans-serif
    run.font.name = 'Arial' 
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    para.space_before = Pt(space_before)
    return para

def add_textbox(slide, text, left, top, width, height,
                size=16, bold=False, color=NAVY, align=PP_ALIGN.LEFT,
                wrap=True, italic=False):
    txb = slide.shapes.add_textbox(left, top, width, height)
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    tf_para(tf, text, size, bold=bold, color=color, align=align)
    para = tf.paragraphs[0]
    if italic and para.runs:
        para.runs[0].font.italic = True
    elif italic:
        run = para.add_run()
        run.text = ""
        run.font.italic = True
    return txb

def header(slide, title_text):
    """Clean, bold header text on every slide."""
    add_textbox(slide, title_text, Inches(0.5), Inches(0.4), Inches(12.0), Inches(0.8),
                size=36, bold=True, color=NAVY)

def footer(slide, text="PolicyForge | Shristi Kumar | Cotiviti Assessment"):
    txb = slide.shapes.add_textbox(Inches(0.5), Inches(7.1), Inches(12), Inches(0.3))
    tf = txb.text_frame
    tf_para(tf, text, 10, color=MID_GRAY)

def bullet_box(slide, items, left, top, width, height, size=15, bullet_color=CYAN, text_color=NAVY):
    txb = slide.shapes.add_textbox(left, top, width, height)
    txb.word_wrap = True
    tf = txb.text_frame
    tf.word_wrap = True
    first = True
    for b_text, text in items:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(8)
        
        # We can simulate bold lead-in by appending runs
        p.text = ""
        r1 = p.add_run()
        r1.text = "■  " + b_text + " "
        r1.font.size = Pt(size)
        r1.font.bold = True
        r1.font.color.rgb = bullet_color
        r1.font.name = 'Arial'

        r2 = p.add_run()
        r2.text = text
        r2.font.size = Pt(size)
        r2.font.color.rgb = text_color
        r2.font.name = 'Arial'


# ─────────────────────────────────────────────────────────────────────────────
# 10 Slides
# ─────────────────────────────────────────────────────────────────────────────

def slide_1_title(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = LIGHT_GRAY

    # Accent block
    add_rect(slide, Inches(0), Inches(0), Inches(0.4), SLIDE_H, CYAN)

    add_textbox(slide, "PolicyForge", Inches(1.0), Inches(2.0), Inches(11), Inches(1.2),
                size=64, bold=True, color=NAVY)
    
    add_textbox(slide, "An Agentic Policy-to-Edit Engine for Medicare Payment Integrity", 
                Inches(1.05), Inches(3.2), Inches(11), Inches(0.8),
                size=24, color=CYAN, bold=True)

    add_textbox(slide, "Shristi Kumar\nCotiviti Intern Assessment\nJuly 2026", 
                Inches(1.05), Inches(5.5), Inches(8), Inches(1.5),
                size=16, color=MID_GRAY)

def slide_2_problem(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE
    header(slide, "The $30B Problem Statement")
    footer(slide)

    # Big metric callout
    add_rect(slide, Inches(0.5), Inches(1.5), Inches(4.5), Inches(5.0), LIGHT_GRAY)
    add_textbox(slide, "$30B+", Inches(0.5), Inches(2.2), Inches(4.5), Inches(1.5),
                size=80, bold=True, color=RED, align=PP_ALIGN.CENTER)
    add_textbox(slide, "Annual Medicare improper payments", Inches(0.5), Inches(3.6), Inches(4.5), Inches(0.8),
                size=16, color=MID_GRAY, align=PP_ALIGN.CENTER)

    # Bullet points
    items = [
        ("The Leakage:", "Rules reside in thousands of unstructured NCDs and LCDs."),
        ("The Bottleneck:", "Manual translation is fundamentally unscalable given CMS velocity."),
        ("The Drain:", "Manual analysis costs ~$56 and takes 45 minutes per policy."),
    ]
    bullet_box(slide, items, Inches(5.5), Inches(2.2), Inches(7.0), Inches(4.0), size=18, bullet_color=NAVY)

def slide_3_solution(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE
    header(slide, "The Solution: PolicyForge Overview")
    footer(slide)

    add_textbox(slide, "Value Proposition: Transitioning from manual text reading to automated, data-driven audit triage.", 
                Inches(0.5), Inches(1.2), Inches(12), Inches(0.5), size=18, bold=True, color=CYAN)

    add_textbox(slide, "An end-to-end, multi-agent engine that reads unstructured Medicare policies, compiles them into executable claim edits, and flags high-utilization provider anomalies in real-world CMS data.", 
                Inches(0.5), Inches(1.8), Inches(12), Inches(0.8), size=16, color=MID_GRAY)

    # Before / After Matrix
    y_start = Inches(3.0)
    add_rect(slide, Inches(1.5), y_start, Inches(4.5), Inches(3.5), LIGHT_GRAY)
    add_textbox(slide, "BEFORE", Inches(1.5), y_start + Inches(0.2), Inches(4.5), Inches(0.5), size=24, bold=True, align=PP_ALIGN.CENTER, color=NAVY)
    bullet_box(slide, [("", "Dense PDF/Text Documents"), ("", "Manual interpretation"), ("", "Static spreadsheets"), ("", "Slow & inconsistent")], 
               Inches(1.8), y_start + Inches(0.8), Inches(4.0), Inches(2.5), size=16, bullet_color=RED)

    add_rect(slide, Inches(7.0), y_start, Inches(4.5), Inches(3.5), CYAN)
    add_textbox(slide, "AFTER", Inches(7.0), y_start + Inches(0.2), Inches(4.5), Inches(0.5), size=24, bold=True, align=PP_ALIGN.CENTER, color=WHITE)
    bullet_box(slide, [("", "Structured JSON objects"), ("", "Python Execution Logic"), ("", "Targeted Provider Flags"), ("", "Fast & reproducible")], 
               Inches(7.3), y_start + Inches(0.8), Inches(4.0), Inches(2.5), size=16, bullet_color=WHITE, text_color=WHITE)

def slide_4_architecture(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = LIGHT_GRAY
    header(slide, "System Architecture & Tech Stack")
    footer(slide)

    nodes = [
        ("Ingestion", "PDF/Text"),
        ("Retriever", "RAG"),
        ("Extractor", "LLM"),
        ("Critic", "Gate"),
        ("Compiler", "Logic"),
        ("Adjudicator", "2σ Outlier"),
        ("Explainer", "Memo")
    ]
    
    start_x = Inches(0.5)
    width = Inches(1.6)
    gap = Inches(0.15)
    y_pos = Inches(2.0)

    for i, (name, sub) in enumerate(nodes):
        x = start_x + i * (width + gap)
        color = CYAN if name in ("Extractor", "Critic", "Compiler") else NAVY
        add_rect(slide, x, y_pos, width, Inches(1.2), color)
        add_textbox(slide, name, x, y_pos + Inches(0.2), width, Inches(0.4), size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_textbox(slide, sub, x, y_pos + Inches(0.6), width, Inches(0.4), size=12, color=WHITE, align=PP_ALIGN.CENTER)
        
        if i < len(nodes) - 1:
            add_textbox(slide, "→", x + width, y_pos + Inches(0.35), gap, Inches(0.5), size=20, bold=True, color=MID_GRAY, align=PP_ALIGN.CENTER)

    # Tech Stack
    add_textbox(slide, "Technology Stack", Inches(0.5), Inches(4.0), Inches(12), Inches(0.5), size=24, bold=True, color=NAVY)
    
    stacks = [
        ("Orchestration:", "LangGraph (Stateful workflow routing)"),
        ("Storage & Vector:", "FAISS, local files"),
        ("Processing & LLM:", "Mistral-large, DuckDB (Analytics), Pandas"),
    ]
    bullet_box(slide, stacks, Inches(0.5), Inches(4.8), Inches(12), Inches(2.0), size=18, bullet_color=CYAN)

def slide_5_orchestration(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE
    header(slide, "LLM Orchestration & The Grounding Gate")
    footer(slide)

    add_textbox(slide, "Why a single prompt fails: Single-pass LLMs suffer from hallucinations and lack validation logic.", 
                Inches(0.5), Inches(1.5), Inches(12), Inches(0.6), size=18, color=RED, bold=True)

    items = [
        ("The Grounding Gate:", "The Critic node mandates verifiable character offsets and source spans before letting a rule compile."),
        ("State Transitions:", "Managed reliably via LangGraph's TypedDict lifecycle."),
    ]
    bullet_box(slide, items, Inches(0.5), Inches(2.3), Inches(6.0), Inches(3.0), size=16, bullet_color=CYAN)

    # Diagram
    x_right = Inches(7.0)
    add_rect(slide, x_right, Inches(2.3), Inches(5.5), Inches(4.0), LIGHT_GRAY)
    
    add_rect(slide, x_right + Inches(1.5), Inches(2.6), Inches(2.5), Inches(0.8), NAVY)
    add_textbox(slide, "Extractor", x_right + Inches(1.5), Inches(2.8), Inches(2.5), Inches(0.4), size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    add_textbox(slide, "↓ proposed rule", x_right + Inches(1.5), Inches(3.4), Inches(2.5), Inches(0.4), size=14, color=MID_GRAY, align=PP_ALIGN.CENTER)

    add_rect(slide, x_right + Inches(1.5), Inches(3.9), Inches(2.5), Inches(0.8), CYAN)
    add_textbox(slide, "Critic Node", x_right + Inches(1.5), Inches(4.1), Inches(2.5), Inches(0.4), size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    add_textbox(slide, "↑ [Fails Verification]\nLoop back for correction", x_right - Inches(0.5), Inches(3.0), Inches(2.5), Inches(1.0), size=12, color=RED, align=PP_ALIGN.CENTER)
    add_textbox(slide, "↓ [Passes]\nTo Compiler", x_right + Inches(1.5), Inches(4.8), Inches(2.5), Inches(1.0), size=12, color=GREEN, align=PP_ALIGN.CENTER)

def slide_6_data(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = LIGHT_GRAY
    header(slide, "Data Engineering & Analytics (DuckDB)")
    footer(slide)

    items = [
        ("Scale:", "Ingestion of CMS Part B dataset containing millions of provider summaries."),
        ("Speed:", "DuckDB dynamically filters targets (e.g., HCPCS 77080/77081) into a local parquet file in ~15 seconds."),
        ("Precision:", "Instead of naive pattern matching, it builds Pydantic schema boundaries for reliable data modeling."),
    ]
    bullet_box(slide, items, Inches(0.5), Inches(1.8), Inches(6.5), Inches(4.0), size=18, bullet_color=CYAN)

    # Code/Data snippet mockup
    x_right = Inches(7.5)
    add_rect(slide, x_right, Inches(1.8), Inches(5.0), Inches(4.0), NAVY)
    add_textbox(slide, "Pydantic Execution Boundary", x_right, Inches(2.0), Inches(5.0), Inches(0.5), size=14, bold=True, color=CYAN, align=PP_ALIGN.CENTER)
    
    code = """class PolicyCriteria(BaseModel):
    policy_id: str
    target_hcpcs_codes: List[str]
    frequency_limit_months: int
    review_required: bool

# DuckDB query generation
SELECT provider_id, SUM(services) 
FROM cms_data
WHERE hcpcs IN (target_hcpcs_codes)
GROUP BY provider_id;"""
    
    add_textbox(slide, code, x_right + Inches(0.3), Inches(2.6), Inches(4.4), Inches(3.0), size=14, color=WHITE)

def slide_7_metrics(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE
    header(slide, "Operational Results & Metrics")
    footer(slide)

    metrics = [
        ("98.2%", "Mean HCPCS F1", "Accuracy across 15 real test policies.", CYAN),
        ("8 Seconds", "Processing Time", "Slashed from 45 mins. ($3.75 vs $56.25 = 15× cost reduction).", GREEN),
        ("1.8%", "Triage Yield", "Isolating the highest-risk 389 providers out of 21,521.", NAVY)
    ]

    for i, (big, title, sub, col) in enumerate(metrics):
        x = Inches(0.5 + i * 4.2)
        add_rect(slide, x, Inches(2.0), Inches(3.8), Inches(4.0), LIGHT_GRAY)
        add_textbox(slide, big, x, Inches(2.5), Inches(3.8), Inches(1.5), size=64, bold=True, color=col, align=PP_ALIGN.CENTER)
        add_textbox(slide, title, x, Inches(4.0), Inches(3.8), Inches(0.5), size=20, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
        add_textbox(slide, sub, x + Inches(0.3), Inches(4.6), Inches(3.2), Inches(1.0), size=14, color=MID_GRAY, align=PP_ALIGN.CENTER)

def slide_8_clinical_risk(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = LIGHT_GRAY
    header(slide, "Managing Clinical Risk & Human-in-the-Loop")
    footer(slide)

    items = [
        ("Honest Variance:", "Tier 2/3 policies hit 100% F1. Tier 1 (Cancer Screening) dipped to 93.3% (e.g. Colorectal missed 2 of 11 codes)."),
        ("The Danger:", "Missing codes can cause false medical claim denials, affecting patient care."),
        ("The Mitigant:", "Deploy PolicyForge as an automated drafting assistant with mandatory human sign-off for high-stakes policies."),
    ]
    bullet_box(slide, items, Inches(0.5), Inches(1.5), Inches(12.0), Inches(2.5), size=16, bullet_color=NAVY)

    # Bar chart simulation
    y_base = Inches(6.5)
    
    # Tier 3
    add_rect(slide, Inches(1.5), Inches(4.5), Inches(2.5), Inches(2.0), GREEN)
    add_textbox(slide, "100%", Inches(1.5), Inches(4.0), Inches(2.5), Inches(0.5), size=24, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
    add_textbox(slide, "Tier 3\nRoutine", Inches(1.5), Inches(6.6), Inches(2.5), Inches(0.8), size=14, color=NAVY, align=PP_ALIGN.CENTER)

    # Tier 2
    add_rect(slide, Inches(5.0), Inches(4.5), Inches(2.5), Inches(2.0), CYAN)
    add_textbox(slide, "100%", Inches(5.0), Inches(4.0), Inches(2.5), Inches(0.5), size=24, bold=True, color=CYAN, align=PP_ALIGN.CENTER)
    add_textbox(slide, "Tier 2\nImportant", Inches(5.0), Inches(6.6), Inches(2.5), Inches(0.8), size=14, color=NAVY, align=PP_ALIGN.CENTER)

    # Tier 1
    add_rect(slide, Inches(8.5), Inches(5.0), Inches(2.5), Inches(1.5), RED) # Shorter
    add_textbox(slide, "93.3%", Inches(8.5), Inches(4.5), Inches(2.5), Inches(0.5), size=24, bold=True, color=RED, align=PP_ALIGN.CENTER)
    add_textbox(slide, "Tier 1\nCritical (Cancer)", Inches(8.5), Inches(6.6), Inches(2.5), Inches(0.8), size=14, color=NAVY, align=PP_ALIGN.CENTER)


def slide_9_recommendations(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE
    header(slide, "Strategic Recommendations for Cotiviti")
    footer(slide)

    # 3-step chevron roadmap
    steps = [
        ("Phase 1: Deploy Now", "Adopt as an internal ingestion assistant to reduce the 45-minute baseline.", CYAN),
        ("Phase 2: Validation", "Implement cross-policy loops comparing edits against public NCCI datasets.", NAVY),
        ("Phase 3: Pilot", "Pilot autonomous claim execution exclusively on low-risk policy tiers.", GREEN),
    ]

    for i, (title, desc, col) in enumerate(steps):
        x = Inches(0.5 + i * 4.1)
        # We can't perfectly draw chevrons via simple add_rect, but we can use the MSO_SHAPE.CHEVRON
        add_chevron(slide, x, Inches(2.5), Inches(4.0), Inches(2.5), col)
        
        # Overlay text
        add_textbox(slide, title, x + Inches(0.3), Inches(2.8), Inches(3.2), Inches(0.5), size=20, bold=True, color=WHITE)
        add_textbox(slide, desc, x + Inches(0.3), Inches(3.4), Inches(3.2), Inches(1.5), size=14, color=WHITE)

def slide_10_conclusion(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = LIGHT_GRAY
    header(slide, "Conclusion & Core Takeaways")
    footer(slide)

    # Bullet points
    items = [
        ("Fast:", "Reduces manual extraction time from 45 minutes to 8 seconds."),
        ("Mathematical:", "Replaces naive filtering with robust 2σ statistical outlier detection."),
        ("Auditable:", "Critic nodes and Pydantic schemas guarantee data integrity for compliance."),
    ]
    bullet_box(slide, items, Inches(0.5), Inches(1.8), Inches(12.0), Inches(2.5), size=20, bullet_color=CYAN)

    # Bold Quote Block
    add_rect(slide, Inches(1.5), Inches(4.5), Inches(10.3), Inches(1.5), NAVY)
    add_rect(slide, Inches(1.5), Inches(4.5), Inches(0.2), Inches(1.5), CYAN) # Left border
    
    quote = "PolicyForge proves that multi-agent orchestration delivers enterprise-grade medical rule extraction that is fast, mathematically rigorous, and auditable."
    add_textbox(slide, quote, Inches(2.0), Inches(4.8), Inches(9.5), Inches(1.0), size=20, bold=True, color=WHITE, italic=True, align=PP_ALIGN.CENTER)

# ─────────────────────────────────────────────────────────────────────────────
# Main
# ─────────────────────────────────────────────────────────────────────────────

def build_presentation():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_1_title(prs)
    slide_2_problem(prs)
    slide_3_solution(prs)
    slide_4_architecture(prs)
    slide_5_orchestration(prs)
    slide_6_data(prs)
    slide_7_metrics(prs)
    slide_8_clinical_risk(prs)
    slide_9_recommendations(prs)
    slide_10_conclusion(prs)

    path = "../PolicyForge_Presentation.pptx"
    prs.save(path)
    print(f"✅ PowerPoint saved: {path}  ({len(prs.slides)} slides)")
    return path


if __name__ == "__main__":
    build_presentation()
